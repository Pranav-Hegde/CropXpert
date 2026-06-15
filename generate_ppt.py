from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

# --- Constants for Formatting ---
HEADING_FONT_SIZE = Pt(16)
CONTENT_FONT_SIZE = Pt(14)

# Margins
MARGIN_LEFT = Inches(0.5)
MARGIN_TOP = Inches(0.5)
MARGIN_RIGHT = Inches(0.5)
MARGIN_BOTTOM = Inches(0.5)
SLIDE_WIDTH = Inches(10)
SLIDE_HEIGHT = Inches(7.5)
TEXT_WIDTH = SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT

# --- Content Definitions ---

TITLE_PAGE_DATA = {
    "title": "CROP RECOMMENDATION SYSTEM",
    "student_details": "Pranav Ramachandra Hegde\n1DS24MCA001\nDepartment of MCA\n2025-2026",
    "guide_details": "Dr. V. Srinivasan\nAssociate Professor\nDept of MCA\nDSCE"
}

# Slide 2: Abstract
ABSTRACT_POINTS = [
    "Agriculture remains the most significant sector of the global economy, particularly in developing nations.",
    "It serves as the primary livelihood for ~58% of the Indian population and constitutes the backbone of the economy.",
    "The sector is currently vulnerable to fluctuating climatic conditions, declining soil productivity, and erratic monsoons.",
    "Traditional practices often fail to optimize crop selection, leading to 'What to plant?' becoming a risky gamble.",
    "This project presents a 'Crop Recommendation System' that integrates Machine Learning with sustainable metrics.",
    "It utilizes a Random Forest Classifier to process 7 key parameters: Nitrogen, Phosphorus, Potassium, Temperature, Humidity, pH, and Rainfall.",
    "The system identifies the most suitable crop from 22 varieties (e.g., Rice, Maize, Chickpea) with >99% accuracy.",
    "A dedicated 'Carbon Footprint Calculator' quantifies emissions from fertilizers, fuel, and electricity in kgCO2e.",
    "Developed using the MERN Stack (MongoDB, Express, React, Node.js) for the web interface and Python (Flask) for analytics.",
    "The architecture is microservices-based, ensuring scalability, maintainability, and high performance.",
    "Bridging the 'Digital Divide' by providing accessible, scientific tools to rural farmers for precision agriculture."
]

# Slide 3: Introduction
INTRO_POINTS = [
    "Agriculture has been the bedrock of civilization since the Neolithic Revolution, providing essential nutrients.",
    "With the global population approaching 8 billion, food security is a paramount concern for all nations.",
    "In India, agriculture is not just an industry but a critical source of employment and GDP contribution.",
    "The 21st century faces an 'Unprecedented Conversion of Challenges': Climate Change and Land Degradation.",
    "Rising global temperatures and extreme weather events (droughts/floods) make output unpredictable.",
    "The 'Productivity Gap': India has large arable land but lower yield-per-hectare compared to global standards.",
    "Farmers often lack access to timely, scientific data, relying instead on intuition or outdated family traditions.",
    "Precision Agriculture (PA) aims to close this gap using IoT, ML, and GPS to manage soil variability.",
    "This project empowers farmers to transition from 'Intuitive Farming' to 'Data-Driven Entrepreneurship'.",
    "By recommending the 'Right Crop for the Right Soil', we aim to increase income and ensure ecological balance."
]

# Slide 4: Problem Statement
PROBLEM_POINTS = [
    "Suboptimal Yields: Incorrect crop selection leads to yields far below the land's potential capacity.",
    "Economic Waste: Significant financial resources (seeds, labor, fertilizers) are lost on failed crop cycles.",
    "Cycle of Rural Debt: Smallholder farmers are vulnerable; a single failure can lead to long-term financial ruin.",
    "Environmental Degradation: Over-application of chemical fertilizers causes soil acidification and water pollution.",
    "Climate Insecurity: Shifts in climatic zones mean crops suitable 10 years ago may no longer be viable today.",
    "Lack of Sustainability Awareness: Farmers cannot measure their Carbon Footprint, hindering 'Green Agriculture'.",
    "The 'Digital Divide': Advanced research exists in labs but fails to reach the farmer in an accessible format.",
    "Complexity of Existing Tools: Current solutions are often too technical or require expensive hardware sensors.",
    "This project addresses these issues by creating a free, easy-to-use, and scientifically robust web platform."
]

# Slide 5: Objectives
OBJECTIVE_POINTS = [
    "1. High-Precision Prediction Engine: Develop a Random Forest model to categorize 22 crops with >99% accuracy.",
    "2. Environmental Quantification: Implement a Carbon Footprint Calculator to translate inputs into CO2 equivalents.",
    "3. Universal Accessibility: Design a responsive User Interface (UI) using React.js that works on all devices.",
    "4. Scalable Architecture: Use a Microservices approach (Node + Flask) to decouple heavy analytics from the UI.",
    "5. Secure & Robust: Implement JWT authentication and BCrypt hashing to protect user data and history.",
    "6. Knowledge Transfer: Include an AI Chatbot ('Agri-Mitra') to answer queries on pests, timing, and care.",
    "7. Economic Integration: Provide an 'Agri-Store' to link scientific advice with the purchase of correct inputs.",
    "8. Digital History: Allow farmers to track their soil health and prediction history over multiple seasons.",
    "9. Speed & Efficiency: Ensure prediction results are generated in under 800ms for a real-time experience."
]

# Slide 6: Literature Survey
LITERATURE_POINTS = [
    "S. Pudumalar et. al. (2021) [1]: Explored Ensemble Techniques; found Random Forest outperforms SVM/Naive Bayes with higher accuracy and less overfitting.",
    "Monali Paul et. al. (2020) [2]: Analyzed Soil Behavior; highlighted the critical non-linear link between Soil pH and Nutrient Bioavailability.",
    "N. Hemageetha et. al. (2022) [3]: IoT & Smart Farming; demonstrated the value of real-time environmental data (Moisture/Temp) for predictions.",
    "T. Nemecek et. al. (2021) [4]: Life Cycle Assessment (LCA); identified Nitrogenous fertilizers as the primary driver of agricultural greenhouse gases.",
    "V. Srinivasan et. al. [5,6]: Advanced Clustering; established frameworks for handling noise and uncertainty in biological datasets.",
    "Comparative Analysis Findings: Most existing apps are 'Static' (rule-based) or lack 'Sustainability' metrics.",
    "Identified Gap: The need for a system that balances 'Productivity' (Yield) with 'Ecology' (Carbon Footprint).",
    "Algorithmic Choice: Random Forest chosen over Neural Networks due to better performance on tabular agricultural data."
]

# Slide 7: Methodology - Spiral Model
METHODOLOGY_POINTS = [
    "The project adopts the 'Spiral Model', ideal for iterative development, risk analysis, and evolving requirements.",
    "Model Selection: Suitable for projects where user feedback must be integrated frequently (e.g., UI adjustments).",
    "Phase 1 - Planning: Defined requirements, gathered 2200-record Kaggle dataset, finalized MERN stack.",
    "Phase 2 - Risk Analysis: Identified potential risks like Model Bias, Latency, and Data Privacy; proposed mitigation strategies.",
    "Phase 3 - Engineering: Implementation of the 4-tier architecture (Frontend, Backend, ML, Database).",
    "   - Frontend: React.js components for Dashboard and Forms.",
    "   - Backend: Node.js/Express Routes for API management.",
    "   - ML Service: Python Flask app for Random Forest Inference.",
    "Phase 4 - Evaluation: Testing against requirements (Unit, Integration, System) and User Acceptance Testing (UAT).",
    "Iterative Refinement: Tuning Hyperparameters (n_estimators=100) and optimizing MongoDB schemas based on test results."
]

# Slide 8: System Architecture (Text + Image)
ARCH_POINTS = [
    "Design Philosophy: Priorities are Modularity, Extensibility, and Aesthetics (Flat Design).",
    "Architectural Style: Distributed Microservices to allow independent scaling.",
    "1. Presentation Layer (Frontend): React.js (v18) SPA.",
    "   - Uses Axios for asynchronous API calls.",
    "2. Application Layer (Backend): Node.js (v16) with Express.js.",
    "   - Handles Routing, Middleware, and Auth logic.",
    "3. Analytical Layer (ML Engine): Python (v3.9) with Flask.",
    "   - Loads the 'pickle' model files.",
    "4. Data Layer (Persistence): MongoDB Atlas (Cloud NoSQL).",
    "   - Stores User Profiles, Logs, and Carbon Reports.",
    "Data Flow: User Input -> React -> Node API -> Flask API -> Model -> Prediction -> Node -> React -> User.",
    "Security Layers: SSL/TLS for transit, BCrypt for passwords, JWT for sessions."
]

# Slide 9: Algorithm - Random Forest
ALGO_POINTS = [
    "Core Engine: Random Forest Classifier, an ensemble learning method constructing multiple Decision Trees.",
    "Why Random Forest?: 1. Handles High Dimensionality. 2. Robust to Outliers. 3. Reduces Variance (Overfitting).",
    "Training Data: 2200 samples covering 22 crops (Rice, Maize, Cotton, Coffee, Jute, etc.).",
    "features (7): N (Nitrogen), P (Phosphorus), K (Potassium), Temperature, Humidity, pH, Rainfall.",
    "Mechanism: 'Bagging' (Bootstrap Aggregation). Each tree votes for a class; the majority vote wins.",
    "Hyperparameters: n_estimators=100 (Number of trees), criterion='gini' (Impurity measure), random_state=42.",
    "Preprocessing Pipeline: 1. MinMaxScaler (0-1) for N,P,K. 2. StandardScaler (Z-score) for Temp/Rain.",
    "Model Persistence: Trained model is serialized using Python's 'Pickle' library for fast loading in Flask.",
    "Performance: Accuracy > 99.0% on Test Set; Precision/Recall/F1-Score consistently high across all classes."
]

# Slide 10: Module 1 - User Authentication
AUTH_POINTS = [
    "Goal: To provide a secure, personalized workspace for every farmer.",
    "Registration: Captures Name, Email, Password. validation ensures unique email addresses.",
    "Password Security: Uses 'bcryptjs' to hash passwords with a 'Salt' (Cost factor 10) before DB storage.",
    "Login Flow:",
    "   1. User enters credentials.",
    "   2. Backend verifies Hash.",
    "   3. Server signs a JWT (JSON Web Token).",
    "Session Management: JWT is stored in client's LocalStorage. It is sent in the 'Authorization' header for requests.",
    "Authorization: Middleware verifies the token signature on protected routes (e.g., /api/predict).",
    "Error Handling: Provides clear feedback for 'Invalid Credentials', 'Server Error', or 'Network Timeout'.",
    "UI Design: Glassmorphism effect on Login forms for a modern, premium aesthetic."
]

# Slide 11: Module 2 - Crop Recommendation
PREDICT_POINTS = [
    "Goal: To analyze soil/weather parameters and predict the optimal crop.",
    "Input Form: User-friendly interface with validation for numeric ranges (e.g., pH 0-14).",
    "Data Transmission: Inputs are packed into a JSON object and sent to the Node.js Relay.",
    "Microservice Handoff: Node.js forwards the request to the Python Flask Service running on Port 5001.",
    "Inference Steps:",
    "   1. Load Model/Scalers.",
    "   2. Transform Input Vector.",
    "   3. Model.predict().",
    "   4. Return Index.",
    "Result Mapping: The predicted numeric index (0-21) is mapped to a String Name (e.g., 'Rice').",
    "Feedback: The frontend displays the Result Name, a high-quality Image, and cultivation advice.",
    "Latency: Optimized to return results in < 800ms to ensure a 'Real-Time' feel."
]

# Slide 12: Module 3 - Carbon Footprint
CARBON_POINTS = [
    "Goal: To promote 'Sustainability by Design' and Net-Zero Agriculture.",
    "Concept: Translates farming activities into Global Warming Potential (GWP) - kgCO2e.",
    "Input Factors:",
    "   1. Fertilizer (Urea - 1.6 kgCO2/kg).",
    "   2. Fuel (Diesel - 2.68 kgCO2/L).",
    "   3. Electricity (Grid).",
    "Calculation Logic: Sum(Input_Quantity * Emission_Factor) for all inputs.",
    "Standards: Emission factors are derived from IPCC (Intergovernmental Panel on Climate Change) guidelines.",
    "Visualization: Uses 'Chart.js' to render dynamic Pie Charts showing the breakdown of emissions.",
    "Interpretation: Color-coded results (Green/Yellow/Red) indicate if the footprint is Eco-friendly or High.",
    "Actionable Insight: Helps farmers realize that optimizing fertilizer use saves money AND the planet."
]

# Slide 13: Module 4 & 5 - Store & Chatbot
EXTRA_MODULE_POINTS = [
    "Agri-Store (E-Commerce): Integrated marketplace to 'Close the Loop' between advice and action.",
    "   - Function: diverse listing of Seeds, Organic Fertilizers, and Farming Tools.",
    "   - Logic: Displays products relevant to the predicted crop (e.g., Rice seeds after Rice prediction).",
    "   - Benefit: Prevents purchase of unverified/wrong inputs; supports local economy.",
    "Agri-Mitra (Chatbot): AI-powered conversational assistant for 24/7 support.",
    "   - NLP Logic: Keyword mapping and basic Natural Language Processing to understand queries.",
    "   - capabilities: Answers questions on Pest Control, Sowing Timing, Irrigation, and Harvesting.",
    "   - Accessibility: Provides text-based support for farmers who prefer chatting over searching.",
    "   - Psych-Social: Reduces isolation by providing an 'Expert in the Pocket' feeling."
]

# Slide 14: Software Testing
TESTING_POINTS = [
    "Methodology: Agile Testing integrated into every sprint of the development lifecycle.",
    "1. Unit Testing (Jest): Verified individual functions like Token Generation, Input Validation, and JSON Parsing.",
    "2. Integration Testing: Verified communication links (React->Node, Node->Flask, Node->MongoDB).",
    "3. System Testing (Black Box): Validated full workflows (Login -> Prediction -> History -> Logout).",
    "Test Cases (Sample):",
    "   - TC01: Standard Login (Pass). | TC02: Invalid Password (Pass). | TC03: Rice Inputs (Pass).",
    "   - TC04: Negative Values (Pass - Blocked). | TC05: Carbon Calc Logic (Pass).",
    "Performance Testing (JMeter): 100 concurrent users. Avg Latency: 142ms. Max RAM: 450MB.",
    "Security Testing: Validated protection against SQL Injection (via Mongoose) and XSS (via React Escaping).",
    "Browser Compatibility: Verified on Chrome, Firefox, Safari, and Mobile Browsers (Responsive)."
]

# Slide 15: Results & Outcomes
RESULTS_POINTS = [
    "System Stability: Achieved Zero-Downtime deployment during stress testing phases.",
    "User Interface: Mobile-First design successfully renders on devices ranging from 320px to 4k resolution.",
    "Model Accuracy: Consistently correctly classified complex test cases (e.g., distinguishing between Coffee and Tea params).",
    "Sustainability Impact: The Carbon Calculator provided actionable insights in 100% of test scenarios.",
    "Response Time: 95th Percentile latency was under 1.2 seconds, meeting rural connectivity requirements.",
    "User Feedback: UAT participants rated the application 4.5/5 for 'Ease of Use' and 'Clarity of Information'."
]

# Slide 16: Conclusion (Separate)
CONCLUSION_POINTS = [
    "Successful Integration: Merged Machine Learning (Random Forest) with Web Technologies (MERN Stack) for a robust solution.",
    "High Accuracy: Achieved >99% classification accuracy on the test dataset, validating the algorithmic approach.",
    "Sustainability Milestone: First-of-its-kind integration of Carbon Footprint calculations in a recommendation system.",
    "Architectural Robustness: Validated the reliability and scalability of the Microservices architecture (Node.js + Flask).",
    "User Empowerment: Provided a tool that replaces intuition with data-driven facts, reducing agricultural risk.",
    "Economic Impact: Potential to significantly reduce input costs (fertilizers) and avoid crop failure losses.",
    "Accessibility: Mobile-first design ensures usability in remote rural areas with basic internet connectivity.",
    "Scalability: The system successfully handled 100 concurrent users during stress testing without latency degradation.",
    "Social Good: Directly contributes to UN Sustainable Development Goals (SDG 2: Zero Hunger, SDG 13: Climate Action).",
    "Final Verdict: The system is a production-ready prototype capable of transforming local agriculture."
]

# Slide 17: Future Scope (Separate)
FUTURE_POINTS = [
    "IoT Sensor Integration: Deploying LoRaWAN soil kits for automated N-P-K data collection to remove manual entry.",
    "Computer Vision Diagnostics: Using CNNs to detect diseases from leaf images via smartphone cameras.",
    "Satellite Monitoring: Integrating Sentinel-2/Landsat multispectral data for regional soil health mapping.",
    "Language Localization: Translating the UI into 12 Indian regional languages (Hindi, Kannada, Tamil, etc.).",
    "Blockchain Certification: Implementing a distributed ledger for tracking 'Organic' produce to boost farmer income.",
    "Government API Linkage: Connecting with the Soil Health Card (SHC) scheme for seamless data import.",
    "Offline Capability: Developing a native Android App with local model inference (TensorFlow Lite).",
    "Marketplace Expansion: Partnering with FPOs (Farmer Producer Organizations) for the Agri-Store logistics.",
    "Drone Technology: Integration for aerial spraying of fertilizers based on model recommendations.",
    "Global Adaptation: Retraining models on African and Southeast Asian soil datasets for international export."
]

# --- Helper Functions ---

def create_full_slide(prs, title, points, image_path=None):
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    
    # 1. Heading Formatting
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = HEADING_FONT_SIZE
    title_shape.text_frame.paragraphs[0].font.bold = True
    # Position: Very top to save space
    title_shape.top = Inches(0.2)
    title_shape.left = MARGIN_LEFT
    title_shape.width = TEXT_WIDTH
    title_shape.height = Inches(0.8) 
    
    # 2. Content Body Formatting
    body_shape = slide.shapes.placeholders[1]
    
    # Check if image needs to be added (Split screen)
    if image_path and os.path.exists(image_path):
        # Text takes left 55%
        body_shape.left = MARGIN_LEFT
        body_shape.top = Inches(1.0)
        body_shape.width = Inches(5.5)
        body_shape.height = SLIDE_HEIGHT - Inches(1.2)
        
        # Image takes right side
        img_left = Inches(6.2)
        img_top = Inches(1.5)
        # We need to make sure the image fits nicely
        slide.shapes.add_picture(image_path, img_left, img_top, width=Inches(3.5))
        
    else:
        # Full page text
        body_shape.left = MARGIN_LEFT
        body_shape.top = Inches(1.0)
        body_shape.width = TEXT_WIDTH
        body_shape.height = SLIDE_HEIGHT - Inches(1.2)
    
    tf = body_shape.text_frame
    tf.clear()
    tf.word_wrap = True
    
    # Add Points
    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.font.size = CONTENT_FONT_SIZE
        p.font.name = "Arial" 
        p.space_before = Pt(8) 
        p.space_after = Pt(8)
        p.level = 0 
        p.alignment = PP_ALIGN.LEFT

def create_title_slide(prs):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = TITLE_PAGE_DATA["title"]
    title.text_frame.paragraphs[0].font.size = Pt(32) 
    
    # Custom text boxes
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(4), Inches(3))
    tf_left = left_box.text_frame
    tf_left.text = TITLE_PAGE_DATA["student_details"]
    for p in tf_left.paragraphs:
        p.font.size = Pt(18)
        
    right_box = slide.shapes.add_textbox(Inches(5.5), Inches(4), Inches(4), Inches(3))
    tf_right = right_box.text_frame
    tf_right.text = TITLE_PAGE_DATA["guide_details"]
    for p in tf_right.paragraphs:
        p.font.size = Pt(18)
        p.alignment = PP_ALIGN.RIGHT

# Slide 18: References
REFERENCE_POINTS = [
    "1. S. Pudumalar et. al. (2021). 'Crop Recommendation System using Ensemble Techniques'.",
    "2. Monali Paul et. al. (2020). 'Analysis of Soil Behavior and Nutrient Bioavailability'.",
    "3. N. Hemageetha et. al. (2022). 'IoT & Smart Farming with Real-time Data'.",
    "4. T. Nemecek et. al. (2021). 'Life Cycle Assessment of Agricultural Greenhouse Gases'.",
    "5. React.js Official Documentation - https://reactjs.org/",
    "6. Scikit-learn Machine Learning Library - https://scikit-learn.org/",
    "7. Flask Web Framework Documentation - https://flask.palletsprojects.com/",
    "8. MongoDB Database Documentation - https://www.mongodb.com/docs/",
    "9. IPCC Guidelines for National Greenhouse Gas Inventories (2019)."
]

def generate_final_ppt(filename):
    prs = Presentation()
    
    # Slide 1
    create_title_slide(prs)
    
    # Slide 2: Abstract
    create_full_slide(prs, "ABSTRACT", ABSTRACT_POINTS)
    
    # Slide 3: Introduction
    create_full_slide(prs, "INTRODUCTION", INTRO_POINTS)
    
    # Slide 4: Problem Statement
    create_full_slide(prs, "PROBLEM STATEMENT", PROBLEM_POINTS)
    
    # Slide 5: Objectives
    create_full_slide(prs, "PROJECT OBJECTIVES", OBJECTIVE_POINTS)
    
    # Slide 6: Literature Survey
    create_full_slide(prs, "LITERATURE SURVEY", LITERATURE_POINTS)
    
    # Slide 7: Methodology
    create_full_slide(prs, "METHODOLOGY: SPIRAL MODEL", METHODOLOGY_POINTS)
    
    # Slide 8: Architecture - WITH IMAGE
    create_full_slide(prs, "SYSTEM ARCHITECTURE", ARCH_POINTS, image_path="images/architecture.png")
    
    # Slide 9: Algorithm
    create_full_slide(prs, "ALGORITHM: RANDOM FOREST", ALGO_POINTS)
    
    # Slide 10: Auth Module
    create_full_slide(prs, "MODULE 1: AUTHENTICATION", AUTH_POINTS)
    
    # Slide 11: Prediction Module
    create_full_slide(prs, "MODULE 2: CROP RECOMMENDATION", PREDICT_POINTS)
    
    # Slide 12: Carbon Module
    create_full_slide(prs, "MODULE 3: CARBON FOOTPRINT", CARBON_POINTS)
    
    # Slide 13: Store & Chat
    create_full_slide(prs, "MODULE 4 & 5: STORE & CHATBOT", EXTRA_MODULE_POINTS)
    
    # Slide 14: Testing
    create_full_slide(prs, "SOFTWARE TESTING & QA", TESTING_POINTS)

    # Slide 15: Outcomes
    create_full_slide(prs, "RESULTS AND OUTCOMES", RESULTS_POINTS)
    
    # Slide 16: Conclusion
    create_full_slide(prs, "CONCLUSION", CONCLUSION_POINTS)

    # Slide 17: Future Scope
    create_full_slide(prs, "FUTURE SCOPE", FUTURE_POINTS)

    # Slide 18: References
    create_full_slide(prs, "REFERENCES", REFERENCE_POINTS)
    
    # Save
    prs.save(filename)
    print(f"Generated {filename}")

if __name__ == "__main__":
    generate_final_ppt("Crop_Recommendation_System_Ref_Slides.pptx")
